import logging
from pptx import Presentation
from PIL import Image
import os
from io import BytesIO

# Set up basic configuration for logging
logging.basicConfig(level=logging.INFO,
                    format='%(asctime)s - %(levelname)s - %(message)s')


class PPTXConverter:
    def __init__(self, input_file):
        self.input_file = input_file
        try:
            self.prs = Presentation(input_file)
            self.output_dir = 'output_images'
            if not os.path.exists(self.output_dir):
                os.makedirs(self.output_dir)
            logging.info("Presentation loaded and output directory is set.")
        except Exception as e:
            logging.error(
                "Failed to load presentation or create output directory.", exc_info=True)
            raise e

    def convert_slides_to_images(self):
        try:
            for i, slide in enumerate(self.prs.slides):
                image_path = os.path.join(
                    self.output_dir, f'slide_{i + 1}.jpeg')
                text_path = os.path.join(self.output_dir, f'slide_{i + 1}.txt')
                self._save_image(slide, image_path, i)
                self._extract_notes(slide, text_path, i)
        except Exception as e:
            logging.error(
                "An error occurred while converting slides.", exc_info=True)
            raise e

    def _save_image(self, slide, image_path, slide_number):
        try:
            slide_image = slide.shapes._spTree.getchildren()[1]
            image_stream = slide_image.image.blob
            image = Image.open(BytesIO(image_stream))
            image.save(image_path, format='JPEG')
            logging.info(f'Slide {slide_number + 1} saved as {image_path}')
        except Exception as e:
            logging.error(
                f"Failed to save image for slide {slide_number + 1}.", exc_info=True)
            raise e

    def _extract_notes(self, slide, text_path, slide_number):
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                with open(text_path, 'w') as text_file:
                    text_file.write(notes)
                logging.info(
                    f'Notes for Slide {slide_number + 1} saved to {text_path}')
            else:
                logging.info(f'No notes for Slide {slide_number + 1}')
        except Exception as e:
            logging.error(
                f"Failed to extract or save notes for slide {slide_number + 1}.", exc_info=True)
            raise e


def main():
    # Example usage with try-except block
    try:
        converter = PPTXConverter('example.pptx')
        converter.convert_slides_to_images()
    except Exception as e:
        logging.critical(
            "Failed to complete the conversion process.", exc_info=True)


if __name__ == "__main__":
    main()
